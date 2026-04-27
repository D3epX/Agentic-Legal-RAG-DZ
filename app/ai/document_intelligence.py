from __future__ import annotations

import io
import uuid
from dataclasses import asdict

from app.ai.cir import CIRBlock, CanonicalIntermediateRepresentation
from app.ai.storage import get_object_storage
from app.services.documents.processor import get_document_processor
from app.services.language import LanguageService


class DocumentIntelligenceService:
    """Build Canonical Intermediate Representation (CIR) from multi-format inputs."""

    def __init__(self) -> None:
        self.processor = get_document_processor()
        self.lang = LanguageService()
        self.storage = get_object_storage()

    def parse_document(
        self,
        *,
        file_uri: str | None,
        raw_text: str | None,
        filename: str | None,
        mime_type: str | None,
        source_language: str | None,
    ) -> CanonicalIntermediateRepresentation:
        pages: list[dict]
        payload_bytes: bytes | None = None

        if raw_text and raw_text.strip():
            pages = self.processor.extract_text(raw_text)
            detected_lang = source_language or self.lang.detect(raw_text[:2500])
        else:
            if not file_uri:
                raise ValueError("Either raw_text or file_uri is required")
            payload_bytes, resolved_mime = self.storage.read_bytes(file_uri)
            mime = (mime_type or resolved_mime or "").lower()
            name = (filename or "").lower()

            if mime.endswith("pdf") or name.endswith(".pdf"):
                pages = self.processor.extract_pdf_text(payload_bytes)
            elif "word" in mime or name.endswith(".docx"):
                pages = self.processor.extract_docx_text(payload_bytes)
            elif "presentation" in mime or name.endswith(".pptx"):
                pages = self._extract_pptx_text(payload_bytes)
            elif mime.startswith("text/") or name.endswith(".txt"):
                pages = self.processor.extract_text(payload_bytes.decode("utf-8", errors="ignore"))
            elif mime.startswith("image/") or name.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp")):
                pages = self._extract_image_text_with_ocr(payload_bytes)
            else:
                raise ValueError(f"Unsupported file type for CIR: mime={mime}")

            detected_lang = source_language or self.lang.detect("\n".join(p["content"] for p in pages)[:2500])

        blocks: list[CIRBlock] = []
        for page_info in pages:
            content = page_info.get("content", "").strip()
            if not content:
                continue
            for para_idx, paragraph in enumerate(self._paragraph_split(content)):
                blocks.append(
                    CIRBlock(
                        block_id=f"blk-{page_info.get('page', 1)}-{para_idx}",
                        block_type=self._classify_block(paragraph),
                        text=paragraph,
                        page=page_info.get("page"),
                        section=self._extract_section_hint(paragraph),
                        metadata={"page": page_info.get("page")},
                    )
                )

        return CanonicalIntermediateRepresentation(
            document_id=str(uuid.uuid4()),
            source_language=detected_lang,
            metadata={
                "filename": filename,
                "mime_type": mime_type,
                "block_count": len(blocks),
                "has_binary_payload": payload_bytes is not None,
            },
            structural_blocks=blocks,
        )

    @staticmethod
    def _paragraph_split(text: str) -> list[str]:
        parts = [p.strip() for p in text.split("\n") if p.strip()]
        if not parts:
            return [text]
        return parts

    @staticmethod
    def _classify_block(paragraph: str) -> str:
        if len(paragraph) <= 90 and paragraph.endswith(":"):
            return "title"
        if "\t" in paragraph:
            return "table"
        return "paragraph"

    @staticmethod
    def _extract_section_hint(paragraph: str) -> str | None:
        lowered = paragraph.lower()
        if lowered.startswith(("section", "chapter", "part", "article", "المادة", "الفصل")):
            return paragraph[:120]
        return None

    @staticmethod
    def _extract_pptx_text(content: bytes) -> list[dict]:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        pages: list[dict] = []
        for idx, slide in enumerate(prs.slides, start=1):
            bits: list[str] = []
            for shape in slide.shapes:
                text_frame = getattr(shape, "text", "")
                if text_frame:
                    bits.append(str(text_frame).strip())
            joined = "\n".join(b for b in bits if b)
            if joined:
                pages.append({"page": idx, "content": joined})
        return pages

    @staticmethod
    def _extract_image_text_with_ocr(content: bytes) -> list[dict]:
        try:
            import pytesseract
            from PIL import Image

            image = Image.open(io.BytesIO(content))
            text = pytesseract.image_to_string(image)
            return [{"page": 1, "content": text}]
        except Exception:
            raise ValueError("OCR dependency not available for image input")


def cir_to_dict(cir: CanonicalIntermediateRepresentation) -> dict:
    return {
        "document_id": cir.document_id,
        "source_language": cir.source_language,
        "metadata": cir.metadata,
        "structural_blocks": [asdict(b) for b in cir.structural_blocks],
        "semantic_chunks": [asdict(c) for c in cir.semantic_chunks],
    }
